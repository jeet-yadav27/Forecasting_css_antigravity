"""
forecasting/dashboard/ppt_export.py
-----------------------------------
Build a PowerPoint briefing from trained forecast results.
Includes all UI plots (when kaleido is available) plus tables.
"""

from __future__ import annotations

import logging
import os
import tempfile
from datetime import datetime

import numpy as np
import pandas as pd

from forecasting.config import FORECAST_HORIZON, N_CV_FOLDS, WARRANTY_MONTHS
from forecasting.dashboard.builder import (
    make_forecast_df,
    make_global_ranking_df,
    make_summary_df,
)

logger = logging.getLogger(__name__)

# Light-blue table palette (matches UI)
_HDR_BLUE = (37, 99, 235)
_ROW_LIGHT = (239, 246, 255)   # #EFF6FF
_ROW_ALT = (191, 219, 254)     # #BFDBFE slight darker


def _set_run(shape, rgb):
    """Set solid fill on a shape (RGB 0–255 tuple)."""
    from pptx.dml.color import RGBColor
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)


def _add_title_bar(slide, prs, text: str, subtitle: str | None = None):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(1.05),
    )
    _set_run(bar, (15, 23, 42))
    bar.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.05),
        prs.slide_width, Inches(0.08),
    )
    _set_run(accent, _HDR_BLUE)
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.68), Inches(12.5), Inches(0.3))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = RGBColor(148, 163, 184)


def _add_table(slide, df: pd.DataFrame, left, top, width, height, max_rows: int = 12,
               highlight_star_col: str | None = "Better"):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    if df is None or df.empty:
        return
    view = df.head(max_rows).copy()
    rows, cols = view.shape[0] + 1, view.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for j, col in enumerate(view.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(10)
            para.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*_HDR_BLUE)

    for i, (_, row) in enumerate(view.iterrows(), start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            if isinstance(val, float) and np.isfinite(val):
                cell.text = f"{val:.2f}" if abs(val) < 1000 else f"{val:,.1f}"
            else:
                cell.text = "" if pd.isna(val) else str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9)
                para.font.color.rgb = RGBColor(15, 23, 42)
            cell.fill.solid()
            star = False
            if highlight_star_col and highlight_star_col in view.columns:
                star = str(view.iloc[i - 1].get(highlight_star_col, "")).strip() in ("★", "*")
            cell.fill.fore_color.rgb = RGBColor(
                *( (187, 247, 208) if star else (_ROW_ALT if i % 2 == 0 else _ROW_LIGHT) )
            )


def _try_chart_png(fig, path: str) -> bool:
    """Export Plotly figure to PNG (requires kaleido)."""
    try:
        fig.write_image(path, width=1280, height=720, scale=1)
        return os.path.exists(path) and os.path.getsize(path) > 0
    except Exception as exc:
        logger.debug("Chart PNG export failed: %s", exc)
        return False


def _add_chart_slide(prs, blank, tmpdir, title: str, fig_fn, *args,
                     subtitle: str | None = None, fname: str = "chart.png") -> bool:
    """Build figure, export PNG, add full-bleed chart slide. Returns success."""
    from pptx.util import Inches

    try:
        fig = fig_fn(*args) if callable(fig_fn) else fig_fn
        if fig is None:
            return False
        png = os.path.join(tmpdir, fname)
        if not _try_chart_png(fig, png):
            return False
        slide = prs.slides.add_slide(blank)
        _add_title_bar(slide, prs, title, subtitle)
        # Tall dashboard figures need a bit less width so height fits
        slide.shapes.add_picture(
            png, Inches(0.45), Inches(1.25), width=Inches(12.4),
        )
        return True
    except Exception as exc:
        logger.debug("Skip chart slide '%s': %s", title, exc)
        return False


def _safe_name(part: str) -> str:
    return "".join(c if c.isalnum() or c in "-_" else "_" for c in str(part))


def build_forecast_pptx(
    results: list[dict],
    *,
    annotated_df: pd.DataFrame | None = None,
    raw: pd.DataFrame | None = None,
    univariate: dict | None = None,
    output_path: str | None = None,
) -> str:
    """
    Create a PowerPoint deck from pipeline results with **all** charts.

    Returns
    -------
    str
        Absolute path to the written .pptx file.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "python-pptx is not installed in this Python environment.\n"
            "Install with:\n"
            "  python -m pip install python-pptx\n"
            "Or activate the project venv first:\n"
            "  .\\venv\\Scripts\\Activate.ps1\n"
            "  pip install -r requirements.txt"
        ) from exc

    from forecasting.dashboard.ppt_charts import (
        accuracy_comparison_df,
        attach_production_from_multivariate,
        executive_rows,
        generate_business_insights,
        holdout_table,
        make_accuracy_comparison_figure,
        make_avp_model_figure,
        make_best_model_forecast_figure,
        make_cpv_cr_dual_figure,
        make_cpv_trend_figure,
        make_cr_trend_figure,
        make_forecast_horizon_figure,
        make_historical_trend_figure,
        make_model_comparison_rmse_figure,
        second_best_model,
        cpv_cr_frames,
    )
    from forecasting.dashboard.builder import (
        make_12m_forecast_figure,
        make_actual_vs_predicted_figure,
        make_annotated_results_figure,
        make_baseline_vs_adjusted_figure,
        make_countermeasure_impact_figure,
        make_cv_fold_figure,
        make_fcok_heatmap,
        make_forecast_df,
        make_historical_claims_figure,
        make_mae_figure,
        make_model_comparison_figure,
        make_odometer_figure,
        make_overview_figure,
        make_part_figure,
        make_production_figure,
        make_vehicle_age_figure,
    )

    uni_map = univariate or {}
    if not results and (annotated_df is None or annotated_df.empty) and not uni_map:
        raise ValueError("No results to export. Train a part or run univariate analysis first.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    charts_ok = 0
    charts_fail = 0

    def _chart(*a, **kw):
        nonlocal charts_ok, charts_fail
        if _add_chart_slide(prs, blank, tmpdir, *a, **kw):
            charts_ok += 1
        else:
            charts_fail += 1

    # Title
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
    )
    _set_run(bg, (15, 23, 42))
    bg.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.55),
        prs.slide_width, Inches(0.12),
    )
    _set_run(accent, _HDR_BLUE)
    accent.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.7))
    tp = t.text_frame.paragraphs[0]
    tp.text = "Automotive Warranty Claims Forecasting"
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(255, 255, 255)

    n_parts = len(results)
    n_uni = len(uni_map)
    s = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.5), Inches(1.2))
    sp = s.text_frame.paragraphs[0]
    sp.text = (
        f"Executive briefing  ·  Multivariate {n_parts} part(s)"
        + (f"  ·  Univariate {n_uni} part(s)" if n_uni else "")
        + f"  ·  {FORECAST_HORIZON}-month horizon  ·  "
        f"{N_CV_FOLDS}-fold walk-forward  ·  {WARRANTY_MONTHS}-mo warranty\n"
        f"Generated {now}"
    )
    sp.font.size = Pt(16)
    sp.font.color.rgb = RGBColor(186, 230, 253)

    tmpdir = tempfile.mkdtemp(prefix="ppt_charts_")

    def _section_slide(title: str, subtitle: str):
        sl = prs.slides.add_slide(blank)
        _set_run(
            sl.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                prs.slide_width, prs.slide_height,
            ),
            (15, 23, 42),
        )
        sl.shapes[-1].line.fill.background()
        tb = sl.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        sb = sl.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8))
        sp2 = sb.text_frame.paragraphs[0]
        sp2.text = subtitle
        sp2.font.size = Pt(16)
        sp2.font.color.rgb = RGBColor(186, 230, 253)

    def _exec_slide(r: dict, mode: str):
        sl = prs.slides.add_slide(blank)
        _add_title_bar(sl, prs, "Executive Summary", f"{mode} · {r.get('part')}")
        rows = executive_rows(r, mode)
        for i, (label, value) in enumerate(rows):
            col, rowi = i % 3, i // 3
            left = Inches(0.45 + col * 4.2)
            top = Inches(1.45 + rowi * 1.7)
            card = sl.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.95), Inches(1.5),
            )
            _set_run(card, _ROW_LIGHT)
            card.line.color.rgb = RGBColor(*_HDR_BLUE)
            lb = sl.shapes.add_textbox(left + Inches(0.15), top + Inches(0.18), Inches(3.65), Inches(0.35))
            lp = lb.text_frame.paragraphs[0]
            lp.text = label
            lp.font.size = Pt(11)
            lp.font.color.rgb = RGBColor(71, 85, 105)
            vb = sl.shapes.add_textbox(left + Inches(0.15), top + Inches(0.52), Inches(3.65), Inches(0.8))
            vp = vb.text_frame.paragraphs[0]
            vp.text = str(value)
            vp.font.size = Pt(16)
            vp.font.bold = True
            vp.font.color.rgb = RGBColor(15, 23, 42)

    def _insight_slide(r: dict, acc_df: pd.DataFrame, mode: str):
        sl = prs.slides.add_slide(blank)
        _add_title_bar(sl, prs, "Forecast Insights", f"{mode} · {r.get('part')}")
        lines = generate_business_insights(r, acc_df, mode)
        box = sl.shapes.add_textbox(Inches(0.55), Inches(1.4), Inches(12.2), Inches(5.6))
        tf = box.text_frame
        tf.word_wrap = True
        if not lines:
            tf.paragraphs[0].text = "No insights generated."
            return
        tf.paragraphs[0].text = "• " + lines[0]
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = "• " + line
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(15, 23, 42)
            p.space_before = Pt(8)

    def _required_pack(r: dict, mode: str):
        nonlocal charts_ok, charts_fail
        part = r.get("part", "part")
        safe = f"{_safe_name(part)}_{mode[:3].lower()}"
        acc_df = accuracy_comparison_df(r)
        best = str(r.get("best_model") or "Model")
        second = second_best_model(r)

        _exec_slide(r, mode)

        _chart(
            f"Best Model Forecast - {best}",
            make_best_model_forecast_figure, r,
            subtitle=f"{mode} · historical actual, last-6 predicted, future forecast, CI",
            fname=f"{safe}_best_fc.png",
        )

        avp = holdout_table(r, model=best, n=6)
        sl = prs.slides.add_slide(blank)
        _add_title_bar(
            sl, prs,
            f"Actual vs Predicted — Last 6 Months ({best})",
            "Month · Actual · Predicted · Difference",
        )
        _add_table(sl, avp, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), max_rows=8)

        _chart(
            f"Actual vs Predicted (Best Model) — {best}",
            make_avp_model_figure, r, best,
            fname=f"{safe}_avp_best.png",
        )
        if second:
            _chart(
                f"Actual vs Predicted (Second Best) — {second}",
                make_avp_model_figure, r, second,
                fname=f"{safe}_avp_2nd.png",
            )

        sl = prs.slides.add_slide(blank)
        _add_title_bar(
            sl, prs, "Accuracy Comparison",
            "Accuracy = 100 − |Σ Actual_6M − Σ Predicted_6M| / Σ Actual_6M × 100",
        )
        _add_table(sl, acc_df, Inches(0.4), Inches(1.35), Inches(12.5), Inches(2.6), max_rows=6)
        png = os.path.join(tmpdir, f"{safe}_acc.png")
        if _try_chart_png(make_accuracy_comparison_figure(acc_df), png):
            sl.shapes.add_picture(png, Inches(1.5), Inches(4.1), width=Inches(10.2))
            charts_ok += 1
        else:
            charts_fail += 1

        _chart("Historical Claims Trend", make_historical_trend_figure, r, fname=f"{safe}_hist.png")
        _chart("Model Comparison Chart", make_model_comparison_rmse_figure, r, fname=f"{safe}_cmp.png")
        _chart("Forecast Horizon Projection", make_forecast_horizon_figure, r, fname=f"{safe}_hz.png")
        _chart("CPV Trend Chart", make_cpv_trend_figure, r, fname=f"{safe}_cpv.png")
        _chart("CR Ratio Trend Chart", make_cr_trend_figure, r, fname=f"{safe}_cr.png")

        sl = prs.slides.add_slide(blank)
        _add_title_bar(
            sl, prs, "CPV & CR Ratio Analysis",
            "CPV = Claims / Production  ·  CR = claims rate (app calculation)",
        )
        h, f = cpv_cr_frames(r)
        if not h.empty:
            show = pd.concat([h.tail(4), f], ignore_index=True) if not f.empty else h.tail(8)
            view = show.copy()
            for c in view.columns:
                if c != "Month":
                    view[c] = pd.to_numeric(view[c], errors="coerce").round(4)
            _add_table(
                sl, view, Inches(0.3), Inches(1.3),
                Inches(12.7), Inches(2.4), max_rows=10,
            )
        png = os.path.join(tmpdir, f"{safe}_dual.png")
        if _try_chart_png(make_cpv_cr_dual_figure(r), png):
            sl.shapes.add_picture(png, Inches(0.6), Inches(3.85), width=Inches(12.1))
            charts_ok += 1
        else:
            charts_fail += 1

        _insight_slide(r, acc_df, mode)
        return acc_df


    if results:
        slide = prs.slides.add_slide(blank)
        _add_title_bar(slide, prs, "Portfolio Overview", "Key performance indicators")

        total_hist = sum(float(np.sum(r["claim_vals"])) for r in results)
        total_fc = sum(
            float(np.sum(r.get("best_forecast", r["ensemble_raw"]))) for r in results
        )
        pct = (total_fc - total_hist) / (total_hist + 1e-9) * 100
        bests = [r.get("best_model", "—") for r in results]
        top = max(set(bests), key=bests.count) if bests else "—"

        cards = [
            ("Historical Claims", f"{int(total_hist):,}"),
            ("Forecast (12 mo)", f"{int(round(total_fc)):,}"),
            ("Trend vs History", f"{pct:+.1f}%"),
            ("Top Model", str(top)),
            ("Parts Analysed", str(n_parts)),
        ]
        for i, (label, value) in enumerate(cards):
            left = Inches(0.4 + i * 2.55)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5),
                Inches(2.35), Inches(1.55),
            )
            _set_run(card, _ROW_LIGHT)
            card.line.color.rgb = RGBColor(*_HDR_BLUE)

            vb = slide.shapes.add_textbox(
                left + Inches(0.12), Inches(1.7), Inches(2.1), Inches(0.7),
            )
            vp = vb.text_frame.paragraphs[0]
            vp.text = value
            vp.font.size = Pt(22)
            vp.font.bold = True
            vp.font.color.rgb = RGBColor(15, 23, 42)

            lb = slide.shapes.add_textbox(
                left + Inches(0.12), Inches(2.5), Inches(2.1), Inches(0.4),
            )
            lp = lb.text_frame.paragraphs[0]
            lp.text = label
            lp.font.size = Pt(11)
            lp.font.color.rgb = RGBColor(71, 85, 105)

        sum_df = make_summary_df(results)
        _add_table(
            slide, sum_df, Inches(0.4), Inches(3.35),
            Inches(12.5), Inches(3.7), max_rows=8,
        )

        # Portfolio-level charts
        _chart(
            "All Parts — Historical & Forecast",
            make_overview_figure, results,
            fname="overview.png",
        )
        _chart(
            "Model MAE Comparison",
            make_mae_figure, results,
            fname="mae.png",
        )

        slide = prs.slides.add_slide(blank)
        _add_title_bar(slide, prs, "Model Ranking", "Average walk-forward metrics across parts")
        rank_df = make_global_ranking_df(results)
        _add_table(
            slide, rank_df, Inches(0.5), Inches(1.5),
            Inches(12.3), Inches(5.4), max_rows=10,
        )

        _section_slide(
            "1. Multivariate Forecasting Results",
            "Best-model forecast, last-6-month accuracy, CPV / CR, and insights",
        )

        for r in results:
            part = r["part"]
            safe = _safe_name(part)
            raw_df = raw if raw is not None else r.get("raw_ref")
            _required_pack(r, "Multivariate")

            slide = prs.slides.add_slide(blank)
            _add_title_bar(
                slide, prs,
                f"Part: {part}",
                f"Best model: {r.get('best_model', '—')}",
            )
            fc_df = make_forecast_df(r)
            _add_table(
                slide, fc_df, Inches(0.4), Inches(1.4),
                Inches(12.5), Inches(5.6), max_rows=12,
            )

            # ── All plots for this part ──────────────────────────────────
            _chart(
                f"Part Dashboard — {part}",
                make_part_figure, r,
                subtitle="Historical + forecast + CM factor",
                fname=f"{safe}_part.png",
            )
            _chart(
                f"12-Month Forecast — {part}",
                make_12m_forecast_figure, r,
                subtitle="Baseline vs CM-adjusted when applicable",
                fname=f"{safe}_fc.png",
            )
            _chart(
                f"Actual vs Predicted — {part}",
                make_actual_vs_predicted_figure, r,
                fname=f"{safe}_avp.png",
            )
            _chart(
                f"Historical Claims — {part}",
                make_historical_claims_figure, r,
                fname=f"{safe}_hist.png",
            )
            _chart(
                f"Production — {part}",
                make_production_figure, r,
                fname=f"{safe}_prod.png",
            )
            _chart(
                f"Rolling CV Performance — {part}",
                make_cv_fold_figure, r,
                fname=f"{safe}_cv.png",
            )
            _chart(
                f"Model Comparison — {part}",
                make_model_comparison_figure, r,
                fname=f"{safe}_cmp.png",
            )

            if raw_df is not None:
                _chart(
                    f"FCO/K × Process Heatmap — {part}",
                    make_fcok_heatmap, raw_df, part,
                    fname=f"{safe}_heat.png",
                )
                _chart(
                    f"Vehicle Age — {part}",
                    make_vehicle_age_figure, raw_df, part,
                    fname=f"{safe}_age.png",
                )
                _chart(
                    f"Odometer — {part}",
                    make_odometer_figure, raw_df, part,
                    fname=f"{safe}_odo.png",
                )

            # Economics table
            months = [str(p) for p in r["monthly"]["period"]]
            econ_df = pd.DataFrame({
                "Month": months,
                "Claims": np.asarray(r["claim_vals"]).ravel(),
                "Production": np.asarray(r["production"]).ravel(),
                "Production_Cost": np.asarray(
                    r.get("production_cost", np.full(len(months), np.nan))
                ).ravel()[:len(months)],
                "CPV": np.asarray(r.get("cpv", np.full(len(months), np.nan))).ravel()[:len(months)],
                "Claim_Ratio": np.asarray(
                    r.get("claim_ratio", np.full(len(months), np.nan))
                ).ravel()[:len(months)],
                "Claim_Ratio_1k": np.asarray(
                    r.get("claim_ratio_per_1k", np.full(len(months), np.nan))
                ).ravel()[:len(months)],
            }).round(4)
            slide = prs.slides.add_slide(blank)
            _add_title_bar(
                slide, prs,
                f"Production · Cost · CPV · Claim Ratio — {part}",
                "User-entered production when provided",
            )
            _add_table(
                slide, econ_df, Inches(0.3), Inches(1.35),
                Inches(12.7), Inches(5.6), max_rows=14,
            )

            # CM-adjusted forecast table + charts
            if r.get("has_cm") and r.get("cm_sim"):
                cms = r["cm_sim"]
                adj_df = pd.DataFrame({
                    "Month": [str(p) for p in r["future_periods"]],
                    "Baseline": np.round(cms["original"], 1),
                    "CM Adjusted": np.round(cms["adjusted"], 1),
                    "Monthly Reduction": np.round(cms["monthly_reduction"], 1),
                    "Cumulative Reduction": np.round(cms["cumulative_reduction"], 1),
                })
                slide = prs.slides.add_slide(blank)
                _add_title_bar(
                    slide, prs,
                    f"Countermeasure Impact — {part}",
                    f"CM date {cms.get('cm_month')} · peak FCO/K {cms.get('peak_fcok')} · "
                    f"reduction {cms.get('reduction_pct', 0):.0f}% "
                    f"({cms.get('reduction_source', 'user')}) · "
                    f"−{cms.get('improvement_pct', 0):.1f}% claims",
                )
                _add_table(
                    slide, adj_df, Inches(0.5), Inches(1.4),
                    Inches(12.3), Inches(5.5), max_rows=14,
                )
                _chart(
                    f"CM-Adjusted Forecast — {part}",
                    make_baseline_vs_adjusted_figure, cms, r["future_periods"],
                    subtitle="Baseline vs post-countermeasure claims",
                    fname=f"{safe}_cm_adj.png",
                )
                _chart(
                    f"CM Reduction Detail — {part}",
                    make_countermeasure_impact_figure, cms,
                    fname=f"{safe}_cm_impact.png",
                )

            fe = r.get("forecast_economics") or {}
            if fe:
                fe_df = pd.DataFrame({
                    "Month": [str(p) for p in r["future_periods"]],
                    "Forecast Claims": np.round(r.get("best_forecast", r["ensemble_raw"]), 1),
                    "Future Production": np.round(fe.get("future_production", []), 1),
                    "Claim Ratio / 1k": np.round(fe.get("forecast_claim_ratio_per_1k", []), 3),
                    "Est. Claim Cost": np.round(fe.get("forecast_claim_cost", []), 2),
                })
                slide = prs.slides.add_slide(blank)
                _add_title_bar(
                    slide, prs,
                    f"Forecast Economics — {part}",
                    f"Avg hist CPV={fe.get('avg_hist_cpv')} · "
                    f"Cost/claim={fe.get('cost_per_claim')}",
                )
                _add_table(
                    slide, fe_df, Inches(0.4), Inches(1.4),
                    Inches(12.5), Inches(5.5), max_rows=14,
                )

    if uni_map:
        _section_slide(
            "2. Univariate Analysis Results",
            "Claims-only models · included because univariate analysis was run",
        )
        mv_by_part = {r.get("part"): r for r in (results or [])}
        for part, uni_res in uni_map.items():
            packed = attach_production_from_multivariate(uni_res, mv_by_part.get(part))
            _required_pack(packed, "Univariate")
            rank = packed.get("ranking")
            if rank is not None and not getattr(rank, "empty", True):
                sl = prs.slides.add_slide(blank)
                _add_title_bar(sl, prs, f"Univariate Ranking — {part}", "Selected models")
                _add_table(sl, rank, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), max_rows=10)

    if annotated_df is not None and not annotated_df.empty:
        slide = prs.slides.add_slide(blank)
        _add_title_bar(
            slide, prs,
            "Annotated Walk-Forward Results",
            "Last-N-month evaluation (Error / Accuracy %)",
        )
        show = annotated_df.copy()
        if "Date" in show.columns:
            show["Date"] = pd.to_datetime(show["Date"]).dt.strftime("%Y-%m")
        _add_table(
            slide, show, Inches(0.3), Inches(1.4),
            Inches(12.7), Inches(5.6), max_rows=10,
        )
        _chart(
            "Annotated Forecast Chart",
            make_annotated_results_figure, annotated_df,
            fname="annotated_ppt.png",
        )

    slide = prs.slides.add_slide(blank)
    _add_title_bar(slide, prs, "Notes & Assumptions", now)
    notes = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5))
    nf = notes.text_frame
    nf.word_wrap = True
    chart_note = (
        f"• Charts embedded: {charts_ok}"
        + (f" · skipped (install kaleido): {charts_fail}" if charts_fail else "")
    )
    lines = [
        f"• Forecast horizon: {FORECAST_HORIZON} months",
        f"• Walk-forward folds: {N_CV_FOLDS}",
        f"• Warranty window: {WARRANTY_MONTHS} months",
        "• Production volumes are user-entered when provided (optional)",
        "• CPV in this deck = Claims / Production; CR = claims rate from the app",
        "• 6-month accuracy uses walk-forward test months: "
        "100 − |Σ Actual − Σ Predicted| / Σ Actual × 100",
        "• Univariate slides are included when Tab 3 analysis has been run",
        "• Countermeasure reduces claims after CM date (production-weighted)",
        "• Models: CNN-LSTM · N-BEATS · Transformer · SARIMA (selection-dependent)",
        chart_note,
        "• Tip: pip install kaleido  — required to embed Plotly charts as images",
    ]
    nf.paragraphs[0].text = lines[0]
    nf.paragraphs[0].font.size = Pt(16)
    nf.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)
    for line in lines[1:]:
        p = nf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(15, 23, 42)
        p.space_before = Pt(8)

    if output_path is None:
        fd, output_path = tempfile.mkstemp(
            suffix=".pptx", prefix="warranty_forecast_",
        )
        os.close(fd)
    prs.save(output_path)
    return output_path
